"""Diagramm-Renderer: prozess, zyklus, matrix, pyramide, timeline."""
import math

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from hwz_tokens import (BLAU_SEQ, FARBEN, arrow, box, farbe, set_paras,
                        skaliert, text_color_for)

GAP = Inches(0.12)


def render_diagram(slide, area, spec, folie_nr):
    typ = (spec or {}).get("typ", "")
    fn = {"prozess": _prozess, "zyklus": _zyklus, "matrix": _matrix,
          "pyramide": _pyramide, "timeline": _timeline}.get(typ)
    if fn is None:
        return [f"Folie {folie_nr}: diagram-typ '{typ}' unbekannt "
                f"(prozess, zyklus, matrix, pyramide, timeline)."]
    schritte = spec.get("schritte") or spec.get("ebenen") or spec.get("punkte") \
        or spec.get("quadranten") or []
    if not schritte:
        return [f"Folie {folie_nr}: diagram '{typ}' ohne Elemente."]
    return fn(slide, area, spec, schritte, folie_nr)


def _step(item):
    if isinstance(item, str):
        return item, ""
    return item.get("t", ""), item.get("sub", "")


def _fill_seq(i, n):
    """Blau-Verlauf dunkel -> hell, erstes Element dunkelblau, letztes hellblau."""
    if n <= 1:
        return FARBEN["blau"]
    seq = BLAU_SEQ if n <= len(BLAU_SEQ) else BLAU_SEQ + ["flaeche"]
    return FARBEN[seq[round(i * (len(seq) - 1) / (n - 1))]]


def _label(shape, titel, sub, t_pt, s_pt, fill, align=PP_ALIGN.CENTER):
    col = text_color_for(fill)
    paras = [{"t": titel, "size": t_pt, "bold": True, "color": col, "align": align, "after": 2}]
    if sub:
        paras.append({"t": sub, "size": s_pt, "color": col, "align": align})
    set_paras(shape.text_frame, paras, anchor=MSO_ANCHOR.MIDDLE, margin=Inches(0.08))


def _prozess(slide, area, spec, schritte, folie_nr):
    warnings = []
    if len(schritte) > 6:
        warnings.append(f"Folie {folie_nr}: Prozess auf 6 Schritte gekürzt.")
        schritte = schritte[:6]
    x, y, w, h = area
    n = len(schritte)
    ch_h = min(Inches(1.5), h)
    cy = y + (h - ch_h) / 2
    cw = (w - GAP * (n - 1)) / n
    t_pt = skaliert(15, n, ab=5)
    for i, item in enumerate(schritte):
        titel, sub = _step(item)
        fill = _fill_seq(i, n)
        sp = box(slide, x + i * (cw + GAP), cy, cw, ch_h,
                 shape=MSO_SHAPE.CHEVRON, fill=fill)
        _label(sp, titel, sub, t_pt, skaliert(11, n, ab=5), fill)
    return warnings


def _zyklus(slide, area, spec, schritte, folie_nr):
    warnings = []
    if len(schritte) > 6:
        warnings.append(f"Folie {folie_nr}: Zyklus auf 6 Schritte gekürzt.")
        schritte = schritte[:6]
    x, y, w, h = area
    n = len(schritte)
    node_r = min(Inches(0.75), h / 6)
    ring = min(w, h) / 2 - node_r
    cx, cy = x + w / 2, y + h / 2
    if spec.get("mitte"):
        m_r = max(ring - node_r - Inches(0.08), Inches(0.4))
        mitte = box(slide, cx - m_r, cy - m_r, 2 * m_r, 2 * m_r,
                    shape=MSO_SHAPE.OVAL, fill=FARBEN["flaeche"])
        _label(mitte, str(spec["mitte"]), "", 18, 12, FARBEN["flaeche"])
    pos = []
    for i in range(n):
        ang = -math.pi / 2 + 2 * math.pi * i / n
        pos.append((cx + ring * math.cos(ang), cy + ring * math.sin(ang)))
    for i, item in enumerate(schritte):
        titel, sub = _step(item)
        px, py = pos[i]
        fill = FARBEN[BLAU_SEQ[i % 2 + 1]]  # blau/mittelblau im Wechsel
        sp = box(slide, px - node_r, py - node_r, 2 * node_r, 2 * node_r,
                 shape=MSO_SHAPE.OVAL, fill=fill)
        _label(sp, titel, sub, skaliert(13, n), skaliert(11, n), fill)
    for i in range(n if n >= 2 else 0):
        (x1, y1), (x2, y2) = pos[i], pos[(i + 1) % n]
        dx, dy = x2 - x1, y2 - y1
        dist = math.hypot(dx, dy) or 1
        pad = node_r + Inches(0.06)
        arrow(slide, x1 + dx / dist * pad, y1 + dy / dist * pad,
              x2 - dx / dist * pad, y2 - dy / dist * pad)
    return warnings


def _matrix(slide, area, spec, quadranten, folie_nr):
    warnings = []
    if len(quadranten) != 4:
        warnings.append(f"Folie {folie_nr}: Matrix braucht genau 4 Quadranten "
                        f"(oben-links, oben-rechts, unten-links, unten-rechts).")
        quadranten = (list(quadranten) + [{}] * 4)[:4]
    x, y, w, h = area
    ax = Inches(0.45)  # Platz für Achsen-Beschriftung links/unten
    gx, gy, gw, gh = x + ax, y, w - ax, h - ax
    qw, qh = (gw - GAP) / 2, (gh - GAP) / 2
    for i, q in enumerate(quadranten):
        col, row = i % 2, i // 2
        fill = FARBEN["hellblau"] if q.get("akzent") else FARBEN["flaeche"]
        sp = box(slide, gx + col * (qw + GAP), gy + row * (qh + GAP), qw, qh,
                 fill=fill, radius=0.05)
        colr = text_color_for(fill)
        paras = [{"t": q.get("titel", ""), "size": 15, "bold": True, "color": colr, "after": 3}]
        if q.get("text"):
            paras.append({"t": q["text"], "size": 12, "color": colr})
        set_paras(sp.text_frame, paras, margin=Inches(0.15))
    if spec.get("x"):
        xb = box(slide, gx, gy + gh + Inches(0.08), gw, Inches(0.32), fill=None)
        set_paras(xb.text_frame, [{"t": spec["x"], "size": 12, "bold": True,
                                   "color": FARBEN["dunkelblau"], "align": PP_ALIGN.CENTER}])
    if spec.get("y"):
        # Rotierte Box: volle Quadranten-Höhe als Breite, zentriert um die Achsenmitte
        yb = box(slide, x + Inches(0.02) - gh / 2 + Inches(0.16),
                 gy + gh / 2 - Inches(0.16), gh, Inches(0.32), fill=None)
        yb.rotation = 270
        set_paras(yb.text_frame, [{"t": spec["y"], "size": 12, "bold": True,
                                   "color": FARBEN["dunkelblau"], "align": PP_ALIGN.CENTER}],
                  wrap=False)
    return warnings


def _pyramide(slide, area, spec, ebenen, folie_nr):
    warnings = []
    if len(ebenen) > 5:
        warnings.append(f"Folie {folie_nr}: Pyramide auf 5 Ebenen gekürzt.")
        ebenen = ebenen[:5]
    x, y, w, h = area
    n = len(ebenen)
    lh = (h - GAP * (n - 1)) / n
    for i, item in enumerate(ebenen):
        titel, sub = _step(item)
        lw = w if n == 1 else w * (0.34 + 0.66 * i / (n - 1))  # oben schmal, unten volle Breite
        fill = _fill_seq(i, n)
        sp = box(slide, x + (w - lw) / 2, y + i * (lh + GAP), lw, lh,
                 fill=fill, radius=0.10)
        _label(sp, titel, sub, skaliert(15, n), skaliert(11, n), fill)
    return warnings


def _timeline(slide, area, spec, punkte, folie_nr):
    warnings = []
    if len(punkte) > 7:
        warnings.append(f"Folie {folie_nr}: Timeline auf 7 Punkte gekürzt.")
        punkte = punkte[:7]
    x, y, w, h = area
    n = len(punkte)
    my = y + h / 2
    pad = Inches(0.4)
    line = box(slide, x + pad, my - Pt(1.25), w - 2 * pad, Pt(2.5),
               shape=MSO_SHAPE.RECTANGLE, fill=FARBEN["mittelblau"])
    dot_r = Inches(0.13)
    lw, lh = min(Inches(2.2), (w - 2 * pad) / n), h / 2 - Inches(0.35)
    for i, item in enumerate(punkte):
        titel, sub = _step(item)
        px = x + pad + (w - 2 * pad) * (i / max(n - 1, 1)) if n > 1 else x + w / 2
        akz = farbe(item.get("akzent") if isinstance(item, dict) else None, "hellblau")
        box(slide, px - dot_r, my - dot_r, 2 * dot_r, 2 * dot_r,
            shape=MSO_SHAPE.OVAL, fill=akz, line=FARBEN["weiss"], line_pt=1.5)
        oben = i % 2 == 0 if n > 4 else True  # ab 5 Punkten alternierend
        ly = my - Inches(0.3) - lh if oben else my + Inches(0.3)
        lx = min(max(px - lw / 2, x), x + w - lw)  # Label innerhalb der Inhaltsfläche halten
        lb = box(slide, lx, ly, lw, lh, fill=None)
        paras = [{"t": titel, "size": skaliert(13, n, ab=6), "bold": True,
                  "color": FARBEN["dunkelblau"], "align": PP_ALIGN.CENTER, "after": 2}]
        if sub:
            paras.append({"t": sub, "size": skaliert(11, n, ab=6),
                          "color": FARBEN["schwarz"], "align": PP_ALIGN.CENTER})
        anchor = MSO_ANCHOR.BOTTOM if oben else MSO_ANCHOR.TOP
        set_paras(lb.text_frame, paras, anchor=anchor, margin=Inches(0.03))
    return warnings
