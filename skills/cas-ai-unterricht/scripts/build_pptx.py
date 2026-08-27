#!/usr/bin/env python3
"""Baut eine HWZ-konforme PPTX aus einem JSON-Folienplan.

Aufruf:
    python3 build_pptx.py plan.json output.pptx [vorlage.potx]

Ohne dritten Parameter wird ../assets/HWZ_Praesentationsvorlage_extern_DE.potx genutzt.
Plan-Format: siehe references/hwz-vorlage.md.
"""
import json
import sys
import tempfile
from pathlib import Path
from zipfile import ZipFile

sys.path.insert(0, str(Path(__file__).resolve().parent))

try:
    from pptx import Presentation
    from pptx.util import Pt
except ImportError:
    sys.exit("python-pptx fehlt. Installieren mit: pip3 install python-pptx")

from hwz_cards import content_area, remove_empty_body, render_cards
from hwz_diagramme import render_diagram

DEFAULT_TEMPLATE = Path(__file__).resolve().parent.parent / "assets" / "HWZ_Praesentationsvorlage_extern_DE.potx"


def potx_to_pptx(potx_path: Path) -> Path:
    """python-pptx kann .potx nicht direkt öffnen: Content-Type umschreiben."""
    with ZipFile(potx_path) as zin:
        data = {n: zin.read(n) for n in zin.namelist()}
    ct = data["[Content_Types].xml"].decode("utf-8").replace(
        "presentationml.template.main", "presentationml.presentation.main"
    )
    data["[Content_Types].xml"] = ct.encode("utf-8")
    tmp = Path(tempfile.mkstemp(suffix=".pptx")[1])
    with ZipFile(tmp, "w") as zout:
        for name, blob in data.items():
            zout.writestr(name, blob)
    return tmp


def norm_para(entry):
    if isinstance(entry, str):
        return entry, 0
    return entry.get("t", ""), int(entry.get("lvl", 0))


def fill_text_frame(tf, paras):
    tf.clear()
    for i, entry in enumerate(paras):
        text, lvl = norm_para(entry)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = lvl


def main():
    if len(sys.argv) < 3:
        sys.exit(__doc__)
    plan_path, out_path = Path(sys.argv[1]), Path(sys.argv[2])
    template = Path(sys.argv[3]) if len(sys.argv) > 3 else DEFAULT_TEMPLATE
    if not template.exists():
        sys.exit(f"Vorlage nicht gefunden: {template}")

    src = potx_to_pptx(template) if template.suffix.lower() == ".potx" else template
    prs = Presentation(str(src))
    # Mitgelieferte Folien der Vorlage entfernen — es zählt nur der Plan.
    xml_slides = prs.slides._sldIdLst
    for sld in list(xml_slides):
        rid = sld.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        prs.part.drop_rel(rid)
        xml_slides.remove(sld)
    layouts = {l.name: l for l in prs.slide_layouts}
    plan = json.loads(plan_path.read_text(encoding="utf-8"))
    warnings = []

    for n, spec in enumerate(plan["slides"], 1):
        layout = layouts.get(spec.get("layout", ""))
        if layout is None:
            warnings.append(f"Folie {n}: Layout '{spec.get('layout')}' unbekannt, nutze 'Inhalt normal'.")
            layout = layouts.get("Inhalt normal") or prs.slide_layouts[5]
        slide = prs.slides.add_slide(layout)
        by_idx = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
        by_type = {}
        for ph in slide.placeholders:
            by_type.setdefault(str(ph.placeholder_format.type), []).append(ph)

        title_ph = next((ph for ph in slide.placeholders
                         if str(ph.placeholder_format.type).startswith(("TITLE", "CENTER_TITLE"))), None)
        title_idx = title_ph.placeholder_format.idx if title_ph is not None else None
        if spec.get("title"):
            if title_ph is not None:
                title_ph.text_frame.text = spec["title"]
            else:
                warnings.append(f"Folie {n}: kein Titel-Platzhalter im Layout '{spec['layout']}'.")
        subtitle_ph = next((ph for ph in slide.placeholders
                            if "SUBTITLE" in str(ph.placeholder_format.type)), None)
        subtitle_idx = subtitle_ph.placeholder_format.idx if subtitle_ph is not None else None
        if spec.get("subtitle"):
            if subtitle_ph is not None:
                subtitle_ph.text_frame.text = spec["subtitle"]
            else:
                warnings.append(f"Folie {n}: kein Untertitel-Platzhalter.")

        if spec.get("bullets"):
            body = next((ph for ph in sorted(slide.placeholders, key=lambda p: p.placeholder_format.idx)
                         if ph.has_text_frame
                         and ph.placeholder_format.idx not in (title_idx, subtitle_idx)
                         and not any(k in str(ph.placeholder_format.type) for k in ("DATE", "FOOTER", "SLIDE_NUMBER", "PICTURE"))), None)
            if body is None:
                warnings.append(f"Folie {n}: kein Body-Platzhalter für bullets.")
            else:
                fill_text_frame(body.text_frame, spec["bullets"])

        for idx_str, paras in (spec.get("texts") or {}).items():
            ph = by_idx.get(int(idx_str))
            if ph is None or not ph.has_text_frame:
                warnings.append(f"Folie {n}: Text-idx {idx_str} nicht gefunden.")
                continue
            fill_text_frame(ph.text_frame, paras if isinstance(paras, list) else [paras])

        for idx_str, img in (spec.get("images") or {}).items():
            ph = by_idx.get(int(idx_str))
            if ph is None:
                warnings.append(f"Folie {n}: Bild-idx {idx_str} nicht gefunden.")
            elif not hasattr(ph, "insert_picture"):
                warnings.append(f"Folie {n}: idx {idx_str} ist kein Bild-Platzhalter.")
            elif not img or not Path(img).is_file():
                warnings.append(f"Folie {n}: Bilddatei fehlt: {img!r}")
            else:
                ph.insert_picture(str(img))

        if spec.get("cards") or spec.get("diagram"):
            if spec.get("bullets"):
                warnings.append(f"Folie {n}: bullets und cards/diagram auf derselben Folie "
                                f"überlagern sich — aufteilen.")
            area = content_area(slide, skip_idx=(title_idx, subtitle_idx))
            if spec.get("cards") and spec.get("diagram"):
                warnings.append(f"Folie {n}: cards UND diagram gesetzt — nur cards gerendert, "
                                f"diagram auf eigene Folie legen.")
            if spec.get("cards"):
                warnings += render_cards(slide, area, spec["cards"], n)
            elif spec.get("diagram"):
                warnings += render_diagram(slide, area, spec["diagram"], n)
            if not spec.get("bullets"):
                remove_empty_body(slide, skip_idx=(title_idx, subtitle_idx))

        if spec.get("notes"):
            slide.notes_slide.notes_text_frame.text = spec["notes"]

    prs.save(str(out_path))
    print(f"OK: {len(plan['slides'])} Folien -> {out_path}")
    for w in warnings:
        print(f"WARNUNG: {w}")
    if warnings:
        sys.exit(2)


if __name__ == "__main__":
    main()
