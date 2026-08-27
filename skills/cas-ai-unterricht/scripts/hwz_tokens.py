"""HWZ-Design-Tokens und Shape-Helfer für Cards und Diagramme.

Nur die Farben aus dem HWZ-Theme verwenden. Textfarbe auf farbigen Flächen
IMMER über text_color_for() bestimmen (programmatisch geprüfter Kontrast).
"""
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

FONT = "Aptos"

FARBEN = {
    "dunkelblau": RGBColor(0x00, 0x14, 0x47),
    "blau": RGBColor(0x00, 0x2F, 0xA8),
    "mittelblau": RGBColor(0x2F, 0x56, 0xB9),
    "hellblau": RGBColor(0x3C, 0x6B, 0xDA),
    "flaeche": RGBColor(0xD8, 0xE1, 0xF7),
    "gruen": RGBColor(0x92, 0xB9, 0x3D),
    "gelb": RGBColor(0xFF, 0xAA, 0x00),
    "rot": RGBColor(0xDA, 0x1E, 0x28),
    "weiss": RGBColor(0xFF, 0xFF, 0xFF),
    "schwarz": RGBColor(0x00, 0x00, 0x00),
}
# Blau-Sequenz für Verläufe in Diagrammen (dunkel -> hell)
BLAU_SEQ = ["dunkelblau", "blau", "mittelblau", "hellblau"]

_HELLE_FLAECHEN = {FARBEN["flaeche"], FARBEN["gelb"], FARBEN["gruen"], FARBEN["weiss"]}


def farbe(name, default="hellblau"):
    return FARBEN.get(name or "", FARBEN[default])


def text_color_for(fill_rgb):
    """Kontrastregel: helle Flächen -> Dunkelblau, dunkle Flächen -> Weiss."""
    return FARBEN["dunkelblau"] if fill_rgb in _HELLE_FLAECHEN else FARBEN["weiss"]


def box(slide, x, y, w, h, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
        fill=None, line=None, line_pt=1.0, radius=0.08):
    """Form ohne Schatten, optional Füllung/Rahmen; radius nur für ROUNDED_RECTANGLE."""
    sp = slide.shapes.add_shape(shape, int(x), int(y), int(w), int(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_pt)
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE and sp.adjustments:
        sp.adjustments[0] = radius
    return sp


def set_paras(tf, paras, anchor=MSO_ANCHOR.TOP, margin=Emu(91440) // 2, wrap=True):
    """Textrahmen befüllen. paras: Liste von Dicts {t,size,bold,color,align,after}."""
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, int(margin))
    first = True
    for spec in paras:
        if not spec.get("t"):
            continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = spec.get("align", PP_ALIGN.LEFT)
        if spec.get("after") is not None:
            p.space_after = Pt(spec["after"])
        run = p.add_run()
        run.text = spec["t"]
        f = run.font
        f.name = FONT
        f.size = Pt(spec.get("size", 12))
        f.bold = bool(spec.get("bold"))
        f.color.rgb = spec.get("color", FARBEN["schwarz"])


def arrow(slide, x1, y1, x2, y2, color=None, width_pt=1.75):
    """Gerader Pfeil (Connector mit Dreiecks-Spitze; Spitze via XML, python-pptx kann das nicht nativ)."""
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, int(x1), int(y1), int(x2), int(y2))
    conn.shadow.inherit = False
    conn.line.color.rgb = color or FARBEN["mittelblau"]
    conn.line.width = Pt(width_pt)
    ln = conn.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)
    return conn


def skaliert(base_pt, n, ab=5, faktor=0.85, minimum=11):
    """Schriftgrösse ab n Elementen sanft verkleinern, nie unter das Lesbarkeits-Minimum."""
    size = base_pt if n < ab else base_pt * faktor
    return max(size, minimum)
