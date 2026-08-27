"""Card-Renderer: Karten-Raster und KPI-Karten in der Inhaltsfläche einer Folie."""
import math

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from hwz_tokens import FARBEN, box, farbe, set_paras, skaliert, text_color_for

# Inhaltsfläche von «Inhalt normal» (Fallback, wenn kein Body-Platzhalter da ist)
DEFAULT_AREA = (Inches(0.56), Inches(1.78), Inches(12.20), Inches(4.84))
GAP = Inches(0.25)
BAR_H = Inches(0.07)
PAD = Inches(0.18)


def content_area(slide, skip_idx=()):
    """Freie Inhaltsfläche: Geometrie des Body-/Object-Platzhalters, sonst Standardfläche."""
    for ph in slide.placeholders:
        t = str(ph.placeholder_format.type)
        if ph.placeholder_format.idx in skip_idx:
            continue
        if t.startswith(("BODY", "OBJECT")) and ph.width > Inches(4):
            return (ph.left, ph.top, ph.width, ph.height)
    return DEFAULT_AREA


def remove_empty_body(slide, skip_idx=()):
    """Leeren Body-Platzhalter entfernen, damit im Bearbeiten-Modus kein «Text hinzufügen» stört."""
    for ph in list(slide.placeholders):
        t = str(ph.placeholder_format.type)
        if (t.startswith(("BODY", "OBJECT")) and ph.placeholder_format.idx not in skip_idx
                and ph.has_text_frame and not ph.text_frame.text.strip()):
            ph._element.getparent().remove(ph._element)


def render_cards(slide, area, spec, folie_nr):
    """Karten-Raster. spec: Liste von Karten oder {"items": [...]}.

    Karte: {"kopf","text"?,"kpi"?,"akzent"?}. Mit "kpi" wird die Karte zur Zahlen-Karte.
    Max. 8 Karten, ab 5 zweizeilig.
    """
    warnings = []
    items = spec.get("items") if isinstance(spec, dict) else spec
    if not items:
        return [f"Folie {folie_nr}: cards ohne Inhalt."]
    items = [i if isinstance(i, dict) else {"kopf": str(i)} for i in items]
    if len(items) > 8:
        warnings.append(f"Folie {folie_nr}: {len(items)} Karten, nur 8 werden gesetzt.")
        items = items[:8]
    x, y, w, h = area
    n = len(items)
    cols = n if n <= 4 else math.ceil(n / 2)
    rows = 1 if n <= 4 else 2
    cw = (w - GAP * (cols - 1)) / cols
    ch = (h - GAP * (rows - 1)) / rows
    # Karten nicht höher als nötig: kappen und den Block vertikal zentrieren
    max_ch = Inches(3.0) if any(i.get("kpi") is not None for i in items) else Inches(2.6)
    if rows == 1 and ch > max_ch:
        ch = max_ch
        y = y + (h - ch) / 2
    kopf_pt = skaliert(16, cols, ab=4)
    text_pt = skaliert(13, cols, ab=4)

    for i, item in enumerate(items):
        col, row = i % cols, i // cols
        cx = x + col * (cw + GAP)
        cy = y + row * (ch + GAP)
        akzent = farbe(item.get("akzent"), "hellblau")
        card = box(slide, cx, cy, cw, ch, fill=FARBEN["flaeche"], radius=0.06)
        # Akzentbalken oben, horizontal um den Eckenradius eingerückt
        inset = Inches(0.12)
        box(slide, cx + inset, cy, cw - 2 * inset, BAR_H, fill=akzent, radius=0.5)
        paras = []
        if item.get("kpi") is not None:
            paras.append({"t": str(item["kpi"]), "size": skaliert(40, rows * 2, ab=3),
                          "bold": True, "color": FARBEN["dunkelblau"],
                          "align": PP_ALIGN.CENTER, "after": 2})
            if item.get("kopf"):
                paras.append({"t": item["kopf"], "size": text_pt, "bold": True,
                              "color": FARBEN["dunkelblau"], "align": PP_ALIGN.CENTER, "after": 2})
            if item.get("text"):
                paras.append({"t": item["text"], "size": skaliert(11.5, cols, ab=4),
                              "color": FARBEN["schwarz"], "align": PP_ALIGN.CENTER})
            anchor = MSO_ANCHOR.MIDDLE
        else:
            if item.get("kopf"):
                paras.append({"t": item["kopf"], "size": kopf_pt, "bold": True,
                              "color": FARBEN["dunkelblau"], "after": 4})
            if item.get("text"):
                paras.append({"t": item["text"], "size": text_pt, "color": FARBEN["schwarz"]})
            anchor = MSO_ANCHOR.TOP
        if not paras:
            warnings.append(f"Folie {folie_nr}: Karte {i + 1} ohne kopf/text/kpi.")
            continue
        set_paras(card.text_frame, paras, anchor=anchor, margin=PAD)
    return warnings
